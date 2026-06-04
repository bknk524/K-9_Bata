from __future__ import annotations
from concurrent.futures import ThreadPoolExecutor, as_completed
from dataclasses import dataclass
import os
import re
import json
import hashlib
import sysconfig
import time
import warnings
from pathlib import Path
from typing import Callable, Iterable, List, Dict, Tuple

import chromadb
import numpy as np
from chromadb.config import Settings as ChromaSettings

from docx import Document as DocxDocument
from pptx import Presentation
import openpyxl
from pypdf import PdfReader

from app.embedding_model import load_embedding_model_name
from app.settings import AppSettings

# このファイルは、文書抽出・索引化・検索・埋め込みデバイス選択をまとめて扱う中核ロジック。
SUPPORTED_EXTS = {".docx", ".pptx", ".xlsx", ".pdf", ".txt"}
FILE_NAME_COLLECTION_SUFFIX = "__file_names"
NPU_ONNX_PROVIDER_PRIORITY = (
    "QNNExecutionProvider",
    "OpenVINOExecutionProvider",
    "DmlExecutionProvider",
    "VitisAIExecutionProvider",
)
KIND_SCORE_WEIGHTS = {
    "content": 1.0,
    "file_name": 0.85,
    "folder_name": 0.55,
}
MAX_EXTRACTED_CHARS = 240_000
MAX_DOCX_PARAGRAPHS = 4_000
MAX_DOCX_TABLE_ROWS = 4_000
MAX_PPTX_SLIDES = 300
MAX_XLSX_ROWS_PER_SHEET = 5_000
MAX_XLSX_CELLS_PER_SHEET = 60_000
MAX_PDF_PAGES = 120
SLOW_EXTRACT_SECONDS = 8.0
_chroma_client_cache: Dict[str, chromadb.PersistentClient] = {}
_collection_cache: Dict[Tuple[str, str], object] = {}
_embedder_cache: Dict[Tuple[str, ...], object] = {}
_optional_dependency_errors: Dict[str, Exception] = {}
_openvino_dll_handles = []


def _configure_openvino_dll_paths() -> None:
    if os.name != "nt":
        return

    candidates = [
        Path(sysconfig.get_paths().get("purelib", "")) / "openvino" / "libs",
        Path(sysconfig.get_paths().get("platlib", "")) / "openvino" / "libs",
    ]
    for candidate in candidates:
        if not candidate.exists():
            continue
        candidate_str = str(candidate)
        path_parts = os.environ.get("PATH", "").split(os.pathsep)
        if candidate_str not in path_parts:
            os.environ["PATH"] = candidate_str + os.pathsep + os.environ.get("PATH", "")
        if hasattr(os, "add_dll_directory"):
            try:
                _openvino_dll_handles.append(os.add_dll_directory(candidate_str))
            except OSError:
                pass


_configure_openvino_dll_paths()

import onnxruntime as ort


@dataclass(frozen=True)
class EmbedderSpec:
    backend: str
    resolved_device: str
    note: str
    torch_device: str | None = None
    onnx_provider: str | None = None
    onnx_model_path: str | None = None
    tokenizer_path: str | None = None


class ModelDependencyError(RuntimeError):
    pass


ProgressCallback = Callable[[int, int, str, str], None]


# パス判定、チャンク分割、差分判定などの共通処理。
def file_sha256(path: str) -> str:
    h = hashlib.sha256()
    with open(path, "rb") as f:
        for chunk in iter(lambda: f.read(1024 * 1024), b""):
            h.update(chunk)
    return h.hexdigest()


def normalize_whitespace(text: str) -> str:
    text = text.replace("\u00a0", " ")
    text = re.sub(r"[ \t]+", " ", text)
    text = re.sub(r"\n{3,}", "\n\n", text)
    return text.strip()


def normalized_extension(path: str) -> str:
    return Path(path).suffix.lower()


def is_temporary_office_file(path: str) -> bool:
    return Path(path).name.startswith("~$")


def is_supported_file(path: str) -> bool:
    return not is_temporary_office_file(path) and normalized_extension(path) in SUPPORTED_EXTS


def should_index_file_name(path: str) -> bool:
    return bool(Path(path).name) and not os.path.isdir(path) and not is_temporary_office_file(path)


def should_index_folder_name(path: str) -> bool:
    return os.path.isdir(path) and bool(Path(path).name)


def should_index_name(path: str) -> bool:
    return should_index_file_name(path) or should_index_folder_name(path)


def chunk_text(text: str, chunk_size: int, overlap: int) -> List[str]:
    text = normalize_whitespace(text)
    if not text:
        return []
    chunks: List[str] = []
    start = 0
    n = len(text)
    while start < n:
        end = min(n, start + chunk_size)
        chunks.append(text[start:end])
        if end == n:
            break
        start = max(0, end - overlap)
    return chunks


def iter_name_indexable_paths(folder: str) -> Iterable[str]:
    folder = os.path.abspath(folder)
    for root, _, files in os.walk(folder):
        if should_index_folder_name(root) and os.path.abspath(root) != folder:
            yield root
        for fn in files:
            path = os.path.join(root, fn)
            if os.path.isfile(path) and should_index_file_name(path):
                yield path


def collect_indexable_paths(folder: str) -> List[str]:
    return sorted(os.path.abspath(path) for path in iter_name_indexable_paths(folder))


def filename_to_embedding_text(path: str) -> str:
    file_name = os.path.basename(path)
    stem = Path(file_name).stem
    ext = Path(file_name).suffix.lower().lstrip(".")
    normalized = re.sub(r"[_\-.]+", " ", stem)
    normalized = re.sub(r"(?<=[a-z0-9])(?=[A-Z])", " ", normalized)
    parts = [file_name, stem, normalized.strip()]
    if ext:
        parts.append(ext)
    unique_parts = []
    for part in parts:
        part = normalize_whitespace(part)
        if part and part not in unique_parts:
            unique_parts.append(part)
    return "\n".join(unique_parts)


def path_signature(path: str) -> str:
    if os.path.isdir(path):
        return hashlib.sha256(f"folder|{Path(path).name}".encode("utf-8")).hexdigest()
    if os.path.isfile(path) and is_supported_file(path):
        return file_sha256(path)
    return hashlib.sha256(f"file_name|{Path(path).name}".encode("utf-8")).hexdigest()


def entry_type(path: str) -> str:
    return "folder" if os.path.isdir(path) else "file"


def manifest_entry_is_complete(entry: Dict | None) -> bool:
    return bool(entry) and "entry_type" in entry and "content_indexed" in entry


def manifest_entry_has_size(entry: Dict | None) -> bool:
    return bool(entry) and "size" in entry


def stat_size(path: str) -> int | None:
    return os.path.getsize(path) if os.path.isfile(path) else None


def resolve_path_signature(path: str, prev: Dict | None, *, content_supported: bool, mtime: float, size: int | None) -> Tuple[str, bool]:
    if not prev:
        sha = path_signature(path)
        return sha, False

    if content_supported:
        if manifest_entry_has_size(prev) and prev.get("mtime") == mtime and prev.get("size") == size:
            return prev["sha256"], True
        sha = path_signature(path)
        return sha, prev.get("sha256") == sha

    sha = path_signature(path)
    return sha, prev.get("sha256") == sha


def weighted_similarity_score(kind: str, distance: float) -> float:
    raw_score = 1.0 - float(distance)
    return raw_score * KIND_SCORE_WEIGHTS.get(kind, 1.0)


def get_available_onnx_providers() -> List[str]:
    try:
        return ort.get_available_providers()
    except Exception:
        return []


def _format_dependency_error_message(package_name: str, exc: Exception) -> str:
    detail = str(exc).strip() or repr(exc)
    hint_lines = [
        f"{package_name} の読み込みに失敗しました。",
        detail,
    ]
    if "DLL load failed" in detail or "アプリケーション制御ポリシー" in detail:
        hint_lines.append(
            "Windows が依存ライブラリのネイティブ DLL / pyd をブロックしている可能性があります。"
            " 今回は pyarrow や torch のようなネイティブ拡張が候補です。"
            " 新しい CPU 用 venv を作り直し、CPU 用 requirements のみを再インストールして確認してください。"
        )
    return " ".join(hint_lines)


def get_auto_tokenizer_class():
    try:
        from transformers import AutoTokenizer

        return AutoTokenizer
    except Exception as exc:
        _optional_dependency_errors["transformers"] = exc
        raise ModelDependencyError(_format_dependency_error_message("transformers", exc)) from exc


def get_sentence_transformer_class():
    try:
        from sentence_transformers import SentenceTransformer

        return SentenceTransformer
    except Exception as exc:
        _optional_dependency_errors["sentence_transformers"] = exc
        raise ModelDependencyError(_format_dependency_error_message("sentence-transformers", exc)) from exc


def find_onnx_model_path(model_name: str) -> Path | None:
    model_path = Path(model_name)
    if not model_path.exists():
        return None

    candidates = [
        model_path / "onnx" / "model.onnx",
        model_path / "model.onnx",
    ]
    for candidate in candidates:
        if candidate.exists():
            return candidate.resolve()
    return None


def find_onnx_tokenizer_path(model_name: str, onnx_model_path: Path) -> Path:
    onnx_dir = onnx_model_path.parent
    if (onnx_dir / "tokenizer.json").exists() or (onnx_dir / "tokenizer_config.json").exists():
        return onnx_dir
    return Path(model_name)


def select_npu_onnx_provider(available_providers: Iterable[str]) -> str | None:
    available = set(available_providers)
    for provider in NPU_ONNX_PROVIDER_PRIORITY:
        if provider in available:
            return provider
    return None


def detect_torch_cuda_available() -> bool:
    try:
        import torch

        return torch.cuda.is_available()
    except Exception:
        return False


def detect_torch_xpu_available() -> bool:
    try:
        import torch

        return hasattr(torch, "xpu") and torch.xpu.is_available()
    except Exception:
        return False


def detect_torch_npu_device() -> str | None:
    try:
        import torch

        if hasattr(torch, "npu") and torch.npu.is_available():
            return "npu"
    except Exception:
        return None
    return None


def resolve_embedder_spec(device_setting: str, model_name: str) -> EmbedderSpec:
    device_setting = (device_setting or "auto").lower()
    if device_setting == "gpu":
        device_setting = "cuda"

    cuda_ok = detect_torch_cuda_available()
    xpu_ok = detect_torch_xpu_available()
    torch_npu_device = detect_torch_npu_device()
    onnx_model_path = find_onnx_model_path(model_name)
    onnx_provider = select_npu_onnx_provider(get_available_onnx_providers()) if onnx_model_path else None

    if device_setting == "cpu":
        return EmbedderSpec(backend="torch", resolved_device="cpu", torch_device="cpu", note="CPU指定")

    if device_setting == "cuda":
        if cuda_ok:
            return EmbedderSpec(backend="torch", resolved_device="gpu", torch_device="cuda", note="GPU(CUDA)指定")
        return EmbedderSpec(
            backend="torch",
            resolved_device="cpu",
            torch_device="cpu",
            note="GPU(CUDA)指定でしたが利用不可のためCPUへフォールバック",
        )

    if device_setting == "xpu":
        if xpu_ok:
            return EmbedderSpec(backend="torch", resolved_device="gpu", torch_device="xpu", note="GPU(Intel XPU)指定")
        return EmbedderSpec(
            backend="torch",
            resolved_device="cpu",
            torch_device="cpu",
            note="GPU(Intel XPU)指定でしたが利用不可のためCPUへフォールバック",
        )

    if device_setting == "npu":
        if torch_npu_device:
            return EmbedderSpec(
                backend="torch",
                resolved_device="npu",
                torch_device=torch_npu_device,
                note="NPU指定（Torchバックエンド）",
            )
        if onnx_provider and onnx_model_path is not None:
            tokenizer_path = find_onnx_tokenizer_path(model_name, onnx_model_path)
            return EmbedderSpec(
                backend="onnx",
                resolved_device="npu",
                onnx_provider=onnx_provider,
                onnx_model_path=str(onnx_model_path),
                tokenizer_path=str(tokenizer_path),
                note=f"NPU指定（ONNX Runtime: {onnx_provider}）",
            )
        return EmbedderSpec(
            backend="torch",
            resolved_device="cpu",
            torch_device="cpu",
            note="NPU指定でしたが利用可能なTorch NPU/ONNX providerが無いためCPUへフォールバック",
        )

    if cuda_ok:
        return EmbedderSpec(backend="torch", resolved_device="gpu", torch_device="cuda", note="AUTO判定: GPU(CUDA)")
    if xpu_ok:
        return EmbedderSpec(backend="torch", resolved_device="gpu", torch_device="xpu", note="AUTO判定: GPU(Intel XPU)")
    if torch_npu_device:
        return EmbedderSpec(
            backend="torch",
            resolved_device="npu",
            torch_device=torch_npu_device,
            note="AUTO判定: NPU(Torch)",
        )
    if onnx_provider and onnx_model_path is not None:
        tokenizer_path = find_onnx_tokenizer_path(model_name, onnx_model_path)
        return EmbedderSpec(
            backend="onnx",
            resolved_device="npu",
            onnx_provider=onnx_provider,
            onnx_model_path=str(onnx_model_path),
            tokenizer_path=str(tokenizer_path),
            note=f"AUTO判定: NPU(ONNX Runtime {onnx_provider})",
        )
    return EmbedderSpec(backend="torch", resolved_device="cpu", torch_device="cpu", note="AUTO判定: CPU")


class OnnxRuntimeEmbedder:
    def __init__(self, model_path: str, tokenizer_path: str, provider: str) -> None:
        self.model_path = model_path
        auto_tokenizer = get_auto_tokenizer_class()
        self.tokenizer = auto_tokenizer.from_pretrained(tokenizer_path)
        provider_config = provider
        self.fixed_input_length: int | None = None
        if provider == "OpenVINOExecutionProvider":
            self.fixed_input_length = 512
            provider_config = (
                provider,
                {
                    "device_type": "NPU",
                    "reshape_input": (
                        f"input_ids[1,{self.fixed_input_length}],"
                        f"attention_mask[1,{self.fixed_input_length}]"
                    ),
                },
            )
        self.session = ort.InferenceSession(
            model_path,
            providers=[provider_config, "CPUExecutionProvider"],
        )
        active_providers = self.session.get_providers()
        if provider != "CPUExecutionProvider" and provider not in active_providers:
            raise ModelDependencyError(
                f"{provider} を要求しましたが ONNX Runtime セッションで有効化されませんでした。"
                " onnxruntime-openvino と openvino のバージョン、NPU ドライバ、DLL 探索パスを確認してください。"
                f" 現在の有効 provider: {', '.join(active_providers)}"
            )
        self.input_names = {item.name for item in self.session.get_inputs()}
        self.output_names = [item.name for item in self.session.get_outputs()]

    def encode(
        self,
        texts,
        normalize_embeddings: bool = True,
        batch_size: int = 32,
        show_progress_bar: bool = False,
    ):
        if isinstance(texts, str):
            texts = [texts]
        if self.fixed_input_length is not None:
            batch_size = 1
        vectors: List[np.ndarray] = []
        for start in range(0, len(texts), batch_size):
            batch = texts[start : start + batch_size]
            tokenizer_kwargs = {
                "padding": True,
                "truncation": True,
                "return_tensors": "np",
            }
            if self.fixed_input_length is not None:
                tokenizer_kwargs.update(
                    {
                        "padding": "max_length",
                        "max_length": self.fixed_input_length,
                    }
                )
            tokens = self.tokenizer(batch, **tokenizer_kwargs)
            feeds = {
                name: np.asarray(value, dtype=np.int64)
                for name, value in tokens.items()
                if name in self.input_names
            }
            outputs = self.session.run(None, feeds)
            output_map = dict(zip(self.output_names, outputs))
            embeddings = output_map.get("sentence_embedding")
            if embeddings is None:
                token_embeddings = output_map["token_embeddings"]
                embeddings = token_embeddings[:, 0, :]
            embeddings = np.asarray(embeddings, dtype=np.float32)
            if normalize_embeddings:
                norms = np.linalg.norm(embeddings, axis=1, keepdims=True)
                embeddings = embeddings / np.clip(norms, 1e-12, None)
            vectors.append(embeddings)

        if not vectors:
            return np.empty((0, 0), dtype=np.float32)
        return np.vstack(vectors)


def append_limited(parts: List[str], text: str, current_chars: int, max_chars: int = MAX_EXTRACTED_CHARS) -> int:
    if current_chars >= max_chars:
        return current_chars
    remaining = max_chars - current_chars
    if len(text) > remaining:
        text = text[:remaining]
    parts.append(text)
    return current_chars + len(text)


def is_extract_limit_reached(current_chars: int) -> bool:
    return current_chars >= MAX_EXTRACTED_CHARS


# Office / PDF / txt から本文を抜き出す抽出処理。
def extract_text_docx(path: str) -> str:
    doc = DocxDocument(path)
    parts: List[str] = []
    current_chars = 0
    for p in doc.paragraphs[:MAX_DOCX_PARAGRAPHS]:
        t = p.text.strip()
        if t:
            current_chars = append_limited(parts, t, current_chars)
            if is_extract_limit_reached(current_chars):
                return "\n".join(parts)
    table_rows = 0
    for table in doc.tables:
        for row in table.rows:
            table_rows += 1
            if table_rows > MAX_DOCX_TABLE_ROWS:
                return "\n".join(parts)
            cells = [c.text.strip() for c in row.cells]
            cells = [c for c in cells if c]
            if cells:
                current_chars = append_limited(parts, " | ".join(cells), current_chars)
                if is_extract_limit_reached(current_chars):
                    return "\n".join(parts)
    return "\n".join(parts)


def extract_text_pptx(path: str) -> str:
    prs = Presentation(path)
    parts: List[str] = []
    current_chars = 0
    for i, slide in enumerate(prs.slides, start=1):
        if i > MAX_PPTX_SLIDES:
            break
        slide_texts: List[str] = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                t = (shape.text or "").strip()
                if t:
                    slide_texts.append(t)
        if slide_texts:
            current_chars = append_limited(parts, f"[Slide {i}]\n" + "\n".join(slide_texts), current_chars)
            if is_extract_limit_reached(current_chars):
                break
    return "\n\n".join(parts)


def extract_text_xlsx(path: str) -> str:
    with warnings.catch_warnings():
        warnings.filterwarnings("ignore", message="DrawingML support is incomplete.*")
        wb = openpyxl.load_workbook(path, data_only=True, read_only=True)
    parts: List[str] = []
    current_chars = 0
    try:
        for ws in wb.worksheets:
            cells_seen = 0
            for row_index, row in enumerate(ws.iter_rows(values_only=True), start=1):
                if row_index > MAX_XLSX_ROWS_PER_SHEET:
                    break
                cells_seen += len(row)
                if cells_seen > MAX_XLSX_CELLS_PER_SHEET:
                    break
                vals = [str(v).strip() for v in row if v is not None and str(v).strip() != ""]
                if vals:
                    current_chars = append_limited(parts, f"[Sheet: {ws.title}] " + " / ".join(vals), current_chars)
                    if is_extract_limit_reached(current_chars):
                        return "\n".join(parts)
        return "\n".join(parts)
    finally:
        wb.close()


def extract_text_pdf(path: str) -> str:
    reader = PdfReader(path)
    parts: List[str] = []
    current_chars = 0
    for i, page in enumerate(reader.pages, start=1):
        if i > MAX_PDF_PAGES:
            break
        try:
            text = page.extract_text() or ""
        except Exception:
            text = ""
        text = text.strip()
        if text:
            current_chars = append_limited(parts, f"[Page {i}]\n{text}", current_chars)
            if is_extract_limit_reached(current_chars):
                break
    return "\n\n".join(parts)


def extract_text_txt(path: str) -> str:
    with open(path, "r", encoding="utf-8", errors="ignore") as f:
        return f.read(MAX_EXTRACTED_CHARS)


def extract_text(path: str) -> str:
    ext = normalized_extension(path)
    if ext == ".docx":
        return extract_text_docx(path)
    if ext == ".pptx":
        return extract_text_pptx(path)
    if ext == ".xlsx":
        return extract_text_xlsx(path)
    if ext == ".pdf":
        return extract_text_pdf(path)
    if ext == ".txt":
        return extract_text_txt(path)
    return ""


def extract_text_with_timing(path: str) -> str:
    start = time.perf_counter()
    print(f"[extract] start {path}", flush=True)
    try:
        return extract_text(path)
    finally:
        elapsed = time.perf_counter() - start
        level = "slow" if elapsed >= SLOW_EXTRACT_SECONDS else "done"
        print(f"[extract] {level} {elapsed:.1f}s {path}", flush=True)


# ChromaDB の client / collection を再利用しながら扱う処理。
def get_chroma_client(chroma_dir: str) -> chromadb.PersistentClient:
    chroma_dir = os.path.abspath(chroma_dir)
    client = _chroma_client_cache.get(chroma_dir)
    if client is None:
        client = chromadb.PersistentClient(
            path=chroma_dir,
            settings=ChromaSettings(anonymized_telemetry=False),
        )
        _chroma_client_cache[chroma_dir] = client
    return client


def get_collection(chroma_dir: str, collection: str):
    chroma_dir = os.path.abspath(chroma_dir)
    key = (chroma_dir, collection)
    col = _collection_cache.get(key)
    if col is None:
        client = get_chroma_client(chroma_dir)
        col = client.get_or_create_collection(name=collection)
        _collection_cache[key] = col
    return col


def file_name_collection_name(collection: str) -> str:
    return f"{collection}{FILE_NAME_COLLECTION_SUFFIX}"


def delete_old_chunks_for_file(col, file_path: str) -> None:
    # 同一file_pathの古いチャンクを消してから入れ直す
    try:
        col.delete(where={"file_path": file_path})
    except Exception:
        pass


def delete_old_chunks_for_paths(cols: Iterable, file_paths: Iterable[str]) -> None:
    unique_paths = list(dict.fromkeys(file_paths))
    if not unique_paths:
        return

    for col in cols:
        try:
            col.delete(where={"file_path": {"$in": unique_paths}})
            continue
        except Exception:
            pass

        for file_path in unique_paths:
            delete_old_chunks_for_file(col, file_path)


def prune_stale_files(manifest: Dict[str, Dict], active_files: Iterable[str]) -> List[str]:
    active_set = set(active_files)
    stale_paths = [file_path for file_path in manifest.keys() if file_path not in active_set]

    for file_path in stale_paths:
        manifest.pop(file_path, None)

    return stale_paths


# CPU / GPU / NPU の判定と、埋め込みモデルの読み出し処理。
def resolve_device(device_setting: str, model_name: str | None = None) -> Tuple[str, str]:
    model_name = model_name or load_embedding_model_name()
    spec = resolve_embedder_spec(device_setting, model_name)
    return spec.resolved_device, spec.note


def get_device_option_statuses(model_name: str | None = None) -> Dict[str, bool]:
    model_name = model_name or load_embedding_model_name()
    return {
        "auto": True,
        "cpu": True,
        "cuda": resolve_embedder_spec("cuda", model_name).resolved_device == "gpu",
        "xpu": resolve_embedder_spec("xpu", model_name).resolved_device == "gpu",
        "npu": resolve_embedder_spec("npu", model_name).resolved_device == "npu",
    }


def get_embedder(model_name: str, spec: EmbedderSpec):
    if spec.backend == "onnx":
        key = ("onnx", spec.onnx_model_path or "", spec.onnx_provider or "")
        if key not in _embedder_cache:
            _embedder_cache[key] = OnnxRuntimeEmbedder(
                model_path=spec.onnx_model_path or "",
                tokenizer_path=spec.tokenizer_path or model_name,
                provider=spec.onnx_provider or "CPUExecutionProvider",
            )
        return _embedder_cache[key]

    key = ("torch", model_name, spec.torch_device or "cpu")
    if key not in _embedder_cache:
        sentence_transformer = get_sentence_transformer_class()
        _embedder_cache[key] = sentence_transformer(model_name, device=spec.torch_device)
    return _embedder_cache[key]


# 前回索引との差分を見るための manifest 読み書き。
def manifest_path(chroma_dir: str) -> Path:
    return Path(chroma_dir) / "manifest.json"


def load_manifest(chroma_dir: str) -> Dict[str, Dict]:
    p = manifest_path(chroma_dir)
    if p.exists():
        return json.loads(p.read_text(encoding="utf-8"))
    return {}


def save_manifest(chroma_dir: str, manifest: Dict[str, Dict]) -> None:
    p = manifest_path(chroma_dir)
    p.parent.mkdir(parents=True, exist_ok=True)
    p.write_text(json.dumps(manifest, ensure_ascii=False, indent=2), encoding="utf-8")


def needs_index_update(settings: AppSettings) -> bool:
    docs_dir = os.path.abspath(settings.docs_dir)
    manifest = load_manifest(settings.chroma_dir)
    current_paths = collect_indexable_paths(docs_dir)

    if set(manifest.keys()) != set(current_paths):
        return True

    try:
        file_name_col = get_collection(settings.chroma_dir, file_name_collection_name(settings.collection))
        if file_name_col.count() != len(current_paths):
            return True
    except Exception:
        return True

    for abs_path in current_paths:
        prev = manifest.get(abs_path)
        if not manifest_entry_is_complete(prev):
            return True
        if entry_type(abs_path) != prev.get("entry_type"):
            return True

        content_supported = os.path.isfile(abs_path) and is_supported_file(abs_path)
        if content_supported:
            if not manifest_entry_has_size(prev):
                return True
            if prev.get("mtime") != os.path.getmtime(abs_path):
                return True
            if prev.get("size") != stat_size(abs_path):
                return True
        elif prev.get("content_indexed"):
            return True

    return False


# 実際の索引化本体。差分判定、本文抽出、埋め込み、DB 更新までをまとめて行う。
def index_folder(
    settings: AppSettings,
    progress_callback: ProgressCallback | None = None,
) -> Tuple[int, int, int, int, str]:
    docs_dir = os.path.abspath(settings.docs_dir)
    os.makedirs(settings.chroma_dir, exist_ok=True)

    col = get_collection(settings.chroma_dir, settings.collection)
    file_name_col = get_collection(settings.chroma_dir, file_name_collection_name(settings.collection))
    model_name = load_embedding_model_name()
    spec = resolve_embedder_spec(settings.device, model_name)
    note = spec.note
    model = None

    manifest = load_manifest(settings.chroma_dir)
    current_paths = collect_indexable_paths(docs_dir)
    stale_paths = prune_stale_files(manifest, current_paths)
    delete_old_chunks_for_paths([col, file_name_col], stale_paths)
    removed_files = len(stale_paths)
    rebuild_file_name_index = file_name_col.count() != len(current_paths)

    to_add_ids: List[str] = []
    to_add_docs: List[str] = []
    to_add_metas: List[Dict] = []
    to_add_embs: List[List[float]] = []
    to_add_file_name_ids: List[str] = []
    to_add_file_name_docs: List[str] = []
    to_add_file_name_metas: List[Dict] = []
    to_add_file_name_embs: List[List[float]] = []

    indexed_paths = 0
    skipped_paths = 0
    extraction_error_paths: List[str] = []
    total_paths = len(current_paths)
    pending_entries: List[Dict] = []
    paths_for_parallel_extract: List[str] = []

    if progress_callback is not None:
        progress_callback(0, total_paths, "", "準備中")

    for index, abs_path in enumerate(current_paths, start=1):
        ext = os.path.splitext(abs_path)[1].lower()
        mtime = os.path.getmtime(abs_path)
        size = stat_size(abs_path)

        prev = manifest.get(abs_path)
        content_supported = os.path.isfile(abs_path) and is_supported_file(abs_path)
        sha, file_is_current = resolve_path_signature(
            abs_path,
            prev,
            content_supported=content_supported,
            mtime=mtime,
            size=size,
        )
        needs_manifest_refresh = (
            not manifest_entry_is_complete(prev)
            or (content_supported and not manifest_entry_has_size(prev))
        )
        if file_is_current and not rebuild_file_name_index and not needs_manifest_refresh:
            skipped_paths += 1
            if progress_callback is not None:
                progress_callback(index, total_paths, abs_path, "スキップ")
            continue

        if file_is_current:
            delete_old_chunks_for_paths([file_name_col], [abs_path])
        else:
            delete_old_chunks_for_paths([col, file_name_col], [abs_path])
            if content_supported:
                paths_for_parallel_extract.append(abs_path)

        pending_entries.append({
            "position": index,
            "path": abs_path,
            "ext": ext,
            "mtime": mtime,
            "size": size,
            "sha": sha,
            "file_is_current": file_is_current,
            "content_supported": content_supported,
        })

    extracted_chunks_by_path: Dict[str, List[str]] = {}
    if paths_for_parallel_extract:
        if len(paths_for_parallel_extract) == 1:
            extract_path = paths_for_parallel_extract[0]
            try:
                if progress_callback is not None:
                    target_position = next(
                        entry["position"] for entry in pending_entries
                        if entry["path"] == extract_path
                    )
                    progress_callback(target_position, total_paths, extract_path, "本文抽出")
                extracted_chunks_by_path[extract_path] = chunk_text(
                    extract_text_with_timing(extract_path),
                    settings.chunk_size,
                    settings.chunk_overlap,
                )
            except Exception:
                extracted_chunks_by_path[extract_path] = []
                extraction_error_paths = [extract_path]
        else:
            progress_position_by_path = {
                entry["path"]: entry["position"]
                for entry in pending_entries
                if entry["path"] in paths_for_parallel_extract
            }

            max_workers = min(len(paths_for_parallel_extract), max(2, min(4, os.cpu_count() or 1)))
            with ThreadPoolExecutor(max_workers=max_workers) as executor:
                future_to_path = {
                    executor.submit(extract_text_with_timing, path): path
                    for path in paths_for_parallel_extract
                }
                completed_extracts = 0
                for future in as_completed(future_to_path):
                    abs_path = future_to_path[future]
                    completed_extracts += 1
                    try:
                        extracted_chunks_by_path[abs_path] = chunk_text(
                            future.result(),
                            settings.chunk_size,
                            settings.chunk_overlap,
                        )
                    except Exception:
                        extracted_chunks_by_path[abs_path] = []
                        extraction_error_paths.append(abs_path)

                    if progress_callback is not None:
                        progress_callback(
                            progress_position_by_path.get(abs_path, completed_extracts),
                            total_paths,
                            abs_path,
                            f"本文抽出 {completed_extracts}/{len(paths_for_parallel_extract)}",
                        )

    extraction_error_set = set(extraction_error_paths)

    for entry in pending_entries:
        abs_path = entry["path"]
        ext = entry["ext"]
        mtime = entry["mtime"]
        size = entry["size"]
        sha = entry["sha"]
        file_is_current = entry["file_is_current"]
        content_supported = entry["content_supported"]
        chunks: List[str] = []

        if content_supported and not file_is_current:
            if abs_path in extraction_error_set:
                content_supported = False
            else:
                chunks = extracted_chunks_by_path.get(abs_path, [])

        file_name_text = filename_to_embedding_text(abs_path)
        texts_to_embed = [file_name_text, *chunks]

        if model is None:
            model = get_embedder(model_name, spec)

        if progress_callback is not None:
            progress_callback(entry["position"], total_paths, abs_path, "埋め込み")
        embs = model.encode(
            texts_to_embed,
            normalize_embeddings=True,
            batch_size=32,
            show_progress_bar=False,
        )

        file_name_emb = embs[0]
        file_name_id = hashlib.sha256(f"{abs_path}|{sha}|file_name".encode("utf-8")).hexdigest()
        to_add_file_name_ids.append(file_name_id)
        to_add_file_name_docs.append(file_name_text)
        to_add_file_name_metas.append({
            "file_path": abs_path,
            "file_ext": ext,
            "file_sha256": sha,
            "kind": "folder_name" if os.path.isdir(abs_path) else "file_name",
            "entry_type": entry_type(abs_path),
            "mtime": mtime,
        })
        to_add_file_name_embs.append(file_name_emb.tolist())

        if content_supported and not file_is_current:
            for i, (chunk, emb) in enumerate(zip(chunks, embs[1:])):
                # ★ID衝突回避：file_path + sha + i をhash化
                cid = hashlib.sha256(f"{abs_path}|{sha}|{i}".encode("utf-8")).hexdigest()
                to_add_ids.append(cid)
                to_add_docs.append(chunk)
                to_add_metas.append({
                    "file_path": abs_path,
                    "file_ext": ext,
                    "file_sha256": sha,
                    "chunk_index": i,
                    "kind": "content",
                    "entry_type": "file",
                    "mtime": mtime,
                })
                to_add_embs.append(emb.tolist())

        manifest[abs_path] = {
            "sha256": sha,
            "mtime": mtime,
            "ext": ext,
            "chunks": len(chunks),
            "size": size,
            "content_indexed": content_supported,
            "entry_type": entry_type(abs_path),
        }

        indexed_paths += 1
        if progress_callback is not None:
            progress_callback(entry["position"], total_paths, abs_path, "完了")

    if to_add_ids:
        col.add(ids=to_add_ids, documents=to_add_docs, metadatas=to_add_metas, embeddings=to_add_embs)
    if to_add_file_name_ids:
        file_name_col.add(
            ids=to_add_file_name_ids,
            documents=to_add_file_name_docs,
            metadatas=to_add_file_name_metas,
            embeddings=to_add_file_name_embs,
        )

    save_manifest(settings.chroma_dir, manifest)
    if extraction_error_paths:
        note = f"{note} / 抽出失敗 {len(extraction_error_paths)} 件は名前のみ索引化"
    if progress_callback is not None:
        progress_callback(total_paths, total_paths, "", "完了")
    return indexed_paths, skipped_paths, len(to_add_ids), removed_files, note


def query_collection(col, q_emb: List[float], n_results: int) -> Tuple[List[str], List[Dict], List[float]]:
    available = col.count()
    if n_results <= 0 or available == 0:
        return [], [], []
    n_results = min(n_results, available)

    res = col.query(
        query_embeddings=[q_emb],
        n_results=n_results,
        include=["documents", "metadatas", "distances"],
    )
    return res["documents"][0], res["metadatas"][0], res["distances"][0]


def merge_hits(
    file_map: Dict[str, Dict],
    docs: List[str],
    metas: List[Dict],
    dists: List[float],
    *,
    default_kind: str,
) -> None:
    for doc, meta, dist in zip(docs, metas, dists):
        fp = meta["file_path"]
        hit_kind = meta.get("kind", default_kind)
        score = weighted_similarity_score(hit_kind, dist)
        entry = file_map.setdefault(fp, {"best_score": score, "hits": []})
        entry["best_score"] = max(entry["best_score"], score)
        snippet = doc
        if hit_kind == "content":
            snippet = (doc[:260] + "…") if len(doc) > 260 else doc
        entry["hits"].append({
            "score": score,
            "chunk_index": meta.get("chunk_index"),
            "file_ext": meta.get("file_ext", ""),
            "kind": hit_kind,
            "entry_type": meta.get("entry_type", "file"),
            "snippet": snippet,
        })


def search(settings: AppSettings, query: str) -> List[Tuple[str, float, List[Dict]]]:
    # 本文検索と名前検索の結果をまとめて、ファイル単位で返す。
    col = get_collection(settings.chroma_dir, settings.collection)
    file_name_col = get_collection(settings.chroma_dir, file_name_collection_name(settings.collection))
    model_name = load_embedding_model_name()
    spec = resolve_embedder_spec(settings.device, model_name)
    model = get_embedder(model_name, spec)

    q_emb = model.encode([query], normalize_embeddings=True)[0].tolist()

    file_map: Dict[str, Dict] = {}
    docs, metas, dists = query_collection(col, q_emb, settings.top_k_chunks)
    merge_hits(file_map, docs, metas, dists, default_kind="content")
    file_name_docs, file_name_metas, file_name_dists = query_collection(
        file_name_col,
        q_emb,
        max(settings.top_k_files * 3, settings.top_k_chunks),
    )
    merge_hits(file_map, file_name_docs, file_name_metas, file_name_dists, default_kind="file_name")

    ranked = sorted(file_map.items(), key=lambda kv: kv[1]["best_score"], reverse=True)[: settings.top_k_files]
    # note はUI側で表示したい場合、settings.deviceから resolve_device して出せます
    return [(fp, data["best_score"], data["hits"]) for fp, data in ranked]
